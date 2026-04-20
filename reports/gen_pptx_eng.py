import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
import os

OUTPUT_PPTX = r"G:\My Drive\Magistrale\2year2semester\AML\Project_AML\reports\Project5_Presentation_Final_EN.pptx"
PLOTS_DIR = r"G:\My Drive\Magistrale\2year2semester\AML\Project_AML\reports"

prs = Presentation()

def add_slide(title, content=None, image_path=None):
    slide = prs.slides.add_slide(prs.slide_layouts[5]) # Title only layout
    
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.bold = True
    
    if content:
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.add_paragraph()
        p.text = content
        p.font.size = Pt(20)
        
    if image_path and os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(1), Inches(3.5), height=Inches(3.5))

# Slide 1: Title
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "High-Performance Semantic Correspondence\nProject 5"
subtitle.text = "Team: Osagie, Lapadula, Pugliese, Bellanca\nCourse: Advanced Machine Learning"

# Slide 2: Introduction
add_slide("1. Introduction & Objective", 
          "Objective:\n"
          "- Establish dense semantic correspondences between image pairs of the same category but with varying appearance and pose.\n\n"
          "Challenges:\n"
          "- Large intra-class variations (shapes, textures).\n"
          "- Geometric distortions (viewpoint, scale).\n"
          "- Lack of vast supervised datasets limits full end-to-end training.\n")

# Slide 3: Framework & Libraries
add_slide("2. Framework & Libraries", 
          "- PyTorch & Torchvision: Core deep learning framework for tensor operations and image augmentation.\n"
          "- DINOv2 (facebookresearch): Pre-trained ViT backbone used to extract dense semantic features without supervision.\n"
          "- PEFT (Parameter-Efficient Fine-Tuning): Conceptual framework implementing LoRA and BitFit for lightweight adaptation.\n"
          "- Gradio: Used for creating an interactive, real-time web UI to test and demonstrate robustness dynamically.\n"
          "- Matplotlib & Seaborn: Real-time academic plotting and logging.")

# Slide 4: Methodology & Architecture
add_slide("3. Methodology: Architecture", 
          "1. Feature Extraction: Images pass through frozen DINOv2 (ViT-B/14) at Layer 11.\n"
          "2. Dense Correlation Map: Cosine similarity is computed between source and target feature grids.\n"
          "3. Adaptive Window Mechanism: Hard-Argmax is unstable. We filter the top-K% similarities dynamically to isolate the relevant neighborhood.\n"
          "4. Temperature-Scaled Soft-Argmax: Computes the sub-pixel center of mass acting as a local regularizer.")

# Slide 5: Hyperparameters Setup
add_slide("4. Hyperparameters & Rationale", 
          "- Learning Rate (1e-4): Found optimal via AdamW. Higher LRs disrupted the pre-trained space; lower LRs stalled.\n"
          "- LoRA Rank (r=16): Sweet spot balancing matrix expressiveness and parameter constraints (only 0.9M params).\n"
          "- Temp (T=0.05): Sharpens Soft-Argmax to prevent blur across background pixels, keeping sub-pixel interpolatability.\n"
          "- Masking (Alpha=0.1): Threshold to suppress matches with extreme uncertainty.")

# Slide 6: Baselines
add_slide("5. Zero-Shot Baselines", 
          "Evaluating off-the-shelf Foundation Models without fine-tuning.\n"
          "- DINOv2 strongly outperforms DINOv3 and SAM, proving its superior object-centric representation.", 
          os.path.join(PLOTS_DIR, "01_Baselines_SPair.png"))

# Slide 7: Layer-wise Semantic
add_slide("6. Internal DINOv2 Dynamics (Layer-wise Analysis)", 
          "- Early layers (L4) focus on low-level textures, yielding poor semantic matches.\n"
          "- Middle layers (L8) begin capturing part-level features.\n"
          "- Deep layers (L11) represent high-level semantic concepts, maximizing PCK.",
          os.path.join(PLOTS_DIR, "09_LayerWise_Analysis.png"))

# Slide 8: PEFT
add_slide("7. Parameter-Efficient Fine-Tuning (SPair-71k)", 
          "Implementation of LoRA and BitFit to specialize DINOv2.\n"
          "- LoRA achieves ~80.5% PCK, doubling the baseline zero-shot performance.\n"
          "- BitFit is extremely competitive despite training only bias terms.",
          os.path.join(PLOTS_DIR, "02_PEFT_Main_Results.png"))

# Slide 9: Parameters
add_slide("8. Parameter Efficiency Trade-off", 
          "- LoRA trains only 1.03% of the network (0.9M parameters).\n"
          "- BitFit trains a negligible 0.09% (82K parameters) yet yields 77.8% PCK, demonstrating incredible bias-driven adaptability.",
          os.path.join(PLOTS_DIR, "08_Parameter_Efficiency.png"))

# Slide 10: Adaptive Window Ablation
add_slide("9. Architecture Ablation: Adaptive Window", 
          "- Removing the Adaptive Window causes LoRA to drop ~3% PCK.\n"
          "- The adaptive boundary mechanism is critical to suppress background noise before computing the expected value.",
          os.path.join(PLOTS_DIR, "03_Ablation_AdaptiveWindow.png"))

# Slide 11: Calibration Ablation
add_slide("10. Temperature Calibration (Soft-Argmax)", 
          "- T=0.01 (blue) behaves like Hard-Argmax, losing sub-pixel smoothness.\n"
          "- T=0.05 (green) is the calibrated sweet-spot for local regularization.\n"
          "- T>=0.50 collapses the prediction into a uniform distribution.",
          os.path.join(PLOTS_DIR, "07_Temperature_Calibration.png"))

# Slide 12: Generalization (OOD PF-Pascal)
add_slide("11. Out-of-Distribution Generalization", 
          "- Testing on PF-Pascal without any re-training.\n"
          "- The finetuned models retain a massive +20% gain over exactly equivalent frozen baselines, proving they learned 'how to match' rather than memorizing domain.",
          os.path.join(PLOTS_DIR, "04_Generalization_PFPascal.png"))

# Slide 13: Geometric Robustness
add_slide("12. Geometric Robustness (Rotation)", 
          "- Models degrade gracefully up to 90 degrees.\n"
          "- DINOv3 shows unique resilience at 45 but structurally shares rotation vulnerability (Absolute Positional Encoding).",
          os.path.join(PLOTS_DIR, "06_Robustness_Rotation.png"))

# Slide 14: Conclusion
add_slide("13. Conclusion & Future Work", 
          "Conclusions:\n"
          "- LoRA + Soft-Argmax is a highly effective pipeline for Semantic Corres.\n"
          "- Temperature scaling is essential for proper coordinate extraction.\n\n"
          "Future Work:\n"
          "- Integration of Inference-Time Augmentation (ITA) to solve rotation.\n"
          "- Exploration of multi-scale feature hierarchies for small scale matching.")

prs.save(OUTPUT_PPTX)
print(f"[INFO] English Presentation generated at: {OUTPUT_PPTX}")
